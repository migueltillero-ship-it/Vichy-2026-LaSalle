from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

# --- Configuración de Identidad Institucional ---
archivo_entrada = 'presentacion_base.pptx'
archivo_salida = 'presentacion_parmentier.pptx'
texto_sustituir = 'La Salle'
texto_nuevo = 'Instituto Gastronómico Parmentier'
nuevo_logo_img = 'logo_parmentier.png'

print(f"Iniciando conversión: {texto_sustituir} -> {texto_nuevo}")

if os.path.exists(archivo_entrada) and os.path.exists(nuevo_logo_img):
    prs = Presentation(archivo_entrada)
    
    for slide in prs.slides:
        formas_a_eliminar = []
        
        for shape in slide.shapes:
            # Reemplazo de menciones textuales
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if texto_sustituir in run.text:
                            run.text = run.text.replace(texto_sustituir, texto_nuevo)
            
            # Reemplazo de logotipos (mantiene posición y tamaño exacto)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide.shapes.add_picture(nuevo_logo_img, shape.left, shape.top, shape.width, shape.height)
                formas_a_eliminar.append(shape)
        
        # Limpieza de elementos de la institución anterior
        for shape in formas_a_eliminar:
            elem = shape._element
            elem.getparent().remove(elem)
            
    prs.save(archivo_salida)
    print(f"ÉXITO: Se ha generado '{archivo_salida}' con la nueva identidad.")
else:
    print("ERROR: Verifique que la presentación y el logo tengan los nombres correctos en la carpeta.")