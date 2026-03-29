from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os

# Configuración de archivos
archivo_entrada = 'presentacion_base.pptx'
archivo_salida = 'presentacion_parmentier.pptx'
palabra_clave_anterior = 'La Salle'
palabra_clave_nueva = 'Instituto Gastronómico Parmentier'
nuevo_emblema = 'logo_parmentier.png'

print("Iniciando la actualización institucional para el Instituto Gastronómico Parmentier...")

if os.path.exists(archivo_entrada) and os.path.exists(nuevo_emblema):
    prs = Presentation(archivo_entrada)
    for slide in prs.slides:
        formas_eliminar = []
        for shape in slide.shapes:
            # Reemplazo de texto
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if palabra_clave_anterior in run.text:
                            run.text = run.text.replace(palabra_clave_anterior, palabra_clave_nueva)
            # Reemplazo de logotipos
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide.shapes.add_picture(nuevo_emblema, shape.left, shape.top, shape.width, shape.height)
                formas_eliminar.append(shape)
        # Eliminación de logos anteriores
        for shape in formas_eliminar:
            elem = shape._element
            elem.getparent().remove(elem)
    prs.save(archivo_salida)
    print("ÉXITO: La presentación ha sido actualizada y guardada correctamente.")
else:
    print(f"ERROR: Verifique que '{archivo_entrada}' y '{nuevo_emblema}' estén en la misma carpeta que este script.")